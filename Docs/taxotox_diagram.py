#!/usr/bin/env python3
"""
taxotox_diagram.py
------------------
Generates Docs/taxotox_flowchart.pptx — an A3-landscape PowerPoint diagram
describing the TaxoTox workflow (taxotox_install.R + app.R).

Requirements:
    pip install python-pptx lxml

Run from anywhere:
    python Docs/taxotox_diagram.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Output path ───────────────────────────────────────────────────────────────
OUT = Path(__file__).parent / "taxotox_flowchart.pptx"

# ── Slide size: A3 landscape (42 × 29.7 cm) ──────────────────────────────────
W_SLIDE = 42.0
H_SLIDE = 29.7

# ── Coordinate system ─────────────────────────────────────────────────────────
# TikZ x increases right, y increases UP; origin = ecotox_dl / user_input row.
# PowerPoint x increases right, y increases DOWN; origin = top-left corner.
SCALE = 0.90        # cm per TikZ cm
OX    = 1.5         # left margin on slide (cm)
OY    = 2.4         # top  margin on slide (cm)  [room for title]
YMAX  = 2.5         # highest TikZ y value used (phantom top of offline box)

def _px(tx):            return OX + tx * SCALE
def _py(ty):            return OY + (YMAX - ty) * SCALE
def _cm(v):             return Cm(v)


# ── Color palette (RGB tuples) ─────────────────────────────────────────────────
def _rgb(r, g, b):      return RGBColor(r, g, b)
def _hex(c):            return f"{int(c[0]):02X}{int(c[1]):02X}{int(c[2]):02X}"

PG_F = _rgb(240, 240, 245);  PG_B = _rgb(130, 130, 160)   # process gray
DB_F = _rgb(235, 225, 255);  DB_B = _rgb(120,  80, 200)    # database purple
OL_F = _rgb(220, 232, 245);  OL_B = _rgb( 91, 130, 185)    # offline blue
ON_F = _rgb(220, 245, 228);  ON_B = _rgb( 60, 160,  90)    # online green
IN_F = _rgb(255, 245, 200);  IN_B = _rgb(210, 160,   0)    # input yellow
AL_F = _rgb(180, 230, 160);  AL_B = _rgb( 60, 150,  40);  AL_D = _rgb( 30, 110,  20)
CR_F = _rgb(255, 215, 155);  CR_B = _rgb(210, 120,  20);  CR_D = _rgb(160,  80,   0)
FI_F = _rgb(180, 215, 250);  FI_B = _rgb( 50, 110, 200);  FI_D = _rgb( 20,  70, 160)
OU_F = _rgb(255, 230, 230);  OU_B = _rgb(200,  60,  60)
RR   = _rgb(200,  40,  40)
RA   = _rgb(210, 140,   0)
RG   = _rgb( 40, 150,  50)
WHITE = _rgb(255, 255, 255)
BLACK = _rgb(  0,   0,   0)
MGRAY = _rgb(100, 100, 100)

# ── XML helpers ────────────────────────────────────────────────────────────────

def _solid_fill(parent, color, alpha_pct=100):
    """Append a solidFill child to *parent* with optional transparency."""
    sf  = etree.SubElement(parent, qn("a:solidFill"))
    clr = etree.SubElement(sf, qn("a:srgbClr"))
    clr.set("val", _hex(color))
    if alpha_pct < 100:
        av = etree.SubElement(clr, qn("a:alpha"))
        av.set("val", str(int(alpha_pct * 1000)))
    return sf


def _clear_fill(spPr):
    for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:pattFill", "a:blipFill"):
        e = spPr.find(qn(tag))
        if e is not None:
            spPr.remove(e)


def _apply_fill(shape, fill_color, alpha_pct=100):
    spPr = shape._element.spPr
    _clear_fill(spPr)
    _solid_fill(spPr, fill_color, alpha_pct)


def _apply_line(shape, color, pt=1.1, dashed=False, alpha_pct=100):
    spPr = shape._element.spPr
    ln = spPr.find(qn("a:ln"))
    if ln is not None:
        spPr.remove(ln)
    ln = etree.SubElement(spPr, qn("a:ln"))
    ln.set("w", str(int(pt * 12700)))
    _solid_fill(ln, color, alpha_pct)
    if dashed:
        pd = etree.SubElement(ln, qn("a:prstDash"))
        pd.set("val", "dash")


def _no_line(shape):
    spPr = shape._element.spPr
    ln = spPr.find(qn("a:ln"))
    if ln is not None:
        spPr.remove(ln)
    ln = etree.SubElement(spPr, qn("a:ln"))
    etree.SubElement(ln, qn("a:noFill"))


def _connector_style(conn, color, pt=1.2, dashed=False):
    """Style a connector: color, width, optional dash, arrowhead at end."""
    spPr = conn._element.spPr
    ln = spPr.find(qn("a:ln"))
    if ln is None:
        ln = etree.SubElement(spPr, qn("a:ln"))
    ln.set("w", str(int(pt * 12700)))
    for tag in ("a:noFill", "a:solidFill"):
        e = ln.find(qn(tag))
        if e is not None:
            ln.remove(e)
    _solid_fill(ln, color)
    if dashed:
        pd = ln.find(qn("a:prstDash"))
        if pd is None:
            pd = etree.SubElement(ln, qn("a:prstDash"))
        pd.set("val", "dash")
    # arrowhead at destination end
    for tag, typ in (("a:tailEnd", "none"), ("a:headEnd", "triangle")):
        el = ln.find(qn(tag))
        if el is None:
            el = etree.SubElement(ln, qn(tag))
        el.set("type", typ)
        if typ == "triangle":
            el.set("w", "med")
            el.set("len", "med")


# ── Text helper ────────────────────────────────────────────────────────────────

def _write_text(shape, paragraphs, base_size=7.5):
    """
    paragraphs: list of items, each either:
      - a str → single run, bold if first paragraph
      - a list of (text, bold, italic, color_or_None) tuples
    """
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left  = _cm(0.06)
    tf.margin_right = _cm(0.06)
    tf.margin_top   = _cm(0.04)
    tf.margin_bottom = _cm(0.04)

    for i, item in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(0)
        p.space_after  = Pt(0)

        if isinstance(item, str):
            runs = [(item, (i == 0), False, None)]
        else:
            runs = item

        for text, bold, italic, color in runs:
            run = p.add_run()
            run.text = text
            run.font.size    = Pt(base_size)
            run.font.bold    = bold
            run.font.italic  = italic
            run.font.name    = "Calibri"
            run.font.color.rgb = color if color else BLACK


# ── Shape factories ────────────────────────────────────────────────────────────
# All positions use TikZ-cm coordinates (centre of node).
# w, h are node dimensions in cm.

MSO_RECT       = 1
MSO_PARA       = 2   # parallelogram
MSO_DIAMOND    = 4
MSO_ROUND_RECT = 5
MSO_CAN        = 13  # cylinder (vertical)


def _add_auto(slide, mso_id, tx, ty, w, h, fill, border,
              border_pt=1.1, dashed=False, fill_alpha=100, radius_adj=None):
    """Add an AutoShape centred at TikZ (tx, ty)."""
    from pptx.util import Emu
    left = _cm(_px(tx) - w / 2)
    top  = _cm(_py(ty) - h / 2)
    sp   = slide.shapes.add_shape(mso_id, left, top, _cm(w), _cm(h))
    _apply_fill(sp, fill, fill_alpha)
    _apply_line(sp, border, border_pt, dashed)
    if radius_adj is not None:
        # adjust rounding — works for ROUNDED_RECTANGLE (adj index 0)
        try:
            avLst = sp._element.spPr.prstGeom.avLst
            adj = avLst.find(qn("a:gd"))
            if adj is not None:
                adj.set("fmla", f"val {int(radius_adj * 100000)}")
        except Exception:
            pass
    sp.text_frame.auto_size = None
    return sp


def add_proc(slide, tx, ty, w, h, fill, border, paragraphs, **kw):
    sp = _add_auto(slide, MSO_ROUND_RECT, tx, ty, w, h,
                   fill, border, border_pt=1.1, radius_adj=0.04, **kw)
    _write_text(sp, paragraphs)
    return sp


def add_db(slide, tx, ty, w, h, paragraphs):
    sp = _add_auto(slide, MSO_CAN, tx, ty, w, h, DB_F, DB_B, border_pt=1.1)
    _write_text(sp, paragraphs, base_size=7.0)
    return sp


def add_diamond(slide, tx, ty, w, h, paragraphs):
    sp = _add_auto(slide, MSO_DIAMOND, tx, ty, w, h, PG_F, PG_B, border_pt=1.1)
    _write_text(sp, paragraphs, base_size=7.0)
    return sp


def add_para(slide, tx, ty, w, h, paragraphs):
    """Parallelogram — used for user-data input node."""
    sp = _add_auto(slide, MSO_PARA, tx, ty, w, h, IN_F, IN_B, border_pt=1.1)
    _write_text(sp, paragraphs)
    return sp


def add_box(slide, x1, y1, x2, y2, fill, border,
            border_pt=1.4, dashed=False, fill_alpha=30):
    """Plain rectangle background container. x1<x2, y1>y2 (TikZ coords)."""
    left = _cm(_px(x1))
    top  = _cm(_py(y1))          # y1 is the TOP edge in TikZ (larger y)
    w    = _cm((x2 - x1) * SCALE)
    h    = _cm((y1 - y2) * SCALE)
    sp   = slide.shapes.add_shape(MSO_RECT, left, top, w, h)
    _apply_fill(sp, fill, fill_alpha)
    _apply_line(sp, border, border_pt, dashed)
    _no_line(sp) if border is None else None
    sp.text_frame.text = ""
    _no_line(sp) if border is None else _apply_line(sp, border, border_pt, dashed)
    return sp


def add_label(slide, tx, ty, text, size=8, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    """Floating text label (no fill, no border)."""
    left = _cm(_px(tx))
    top  = _cm(_py(ty))
    tb   = slide.shapes.add_textbox(left, top, _cm(6), _cm(1.0))
    _no_line(tb)
    tf = tb.text_frame
    tf.word_wrap = False
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.name  = "Calibri"
    run.font.color.rgb = color
    return tb


def add_textbox(slide, tx, ty, w, h, paragraphs, fill=WHITE, border=PG_B,
                border_pt=1.0, base_size=7.5, align=PP_ALIGN.LEFT):
    """General-purpose text box with fill and border."""
    left = _cm(_px(tx) - w / 2)
    top  = _cm(_py(ty) - h / 2)
    sp   = slide.shapes.add_shape(MSO_ROUND_RECT, left, top, _cm(w), _cm(h))
    _apply_fill(sp, fill)
    _apply_line(sp, border, border_pt)
    try:
        avLst = sp._element.spPr.prstGeom.avLst
        adj   = avLst.find(qn("a:gd"))
        if adj is not None:
            adj.set("fmla", "val 3000")
    except Exception:
        pass
    tf = sp.text_frame
    tf.word_wrap = True
    tf.margin_left  = _cm(0.12)
    tf.margin_right = _cm(0.12)
    tf.margin_top   = _cm(0.08)
    tf.margin_bottom = _cm(0.08)
    for i, item in enumerate(paragraphs):
        p   = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment  = align
        p.space_before = Pt(0)
        p.space_after  = Pt(1)
        if isinstance(item, str):
            runs = [(item, (i == 0), False, None)]
        else:
            runs = item
        for text, bold, italic, color in runs:
            run = p.add_run()
            run.text = text
            run.font.size   = Pt(base_size)
            run.font.bold   = bold
            run.font.italic = italic
            run.font.name   = "Calibri"
            run.font.color.rgb = color if color else BLACK
    return sp


def arrow(slide, x1, y1, x2, y2, color, pt=1.2, dashed=False):
    """Draw a straight connector with arrowhead at (x2, y2)."""
    from pptx.enum.shapes import MSO_CONNECTOR_TYPE
    bx = _cm(_px(x1));  by = _cm(_py(y1))
    ex = _cm(_px(x2));  ey = _cm(_py(y2))
    conn = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, bx, by, ex, ey)
    _connector_style(conn, color, pt, dashed)
    return conn


def arrow_label(slide, tx, ty, text, color=BLACK, size=6.5):
    """Small label placed near an arrow midpoint."""
    tb = slide.shapes.add_textbox(_cm(_px(tx)), _cm(_py(ty)), _cm(1.8), _cm(0.5))
    _no_line(tb)
    tf = tb.text_frame
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.name  = "Calibri"
    run.font.color.rgb = color
    return tb


# =============================================================================
# NODE POSITIONS  (TikZ-cm coordinates, all absolute)
# =============================================================================
# Derived by tracing the TikZ relative-positioning chain.
# Format:  TX, TY, W, H  (centre x, centre y, width, height) all in cm.

N = {
    # ── Offline chain (x ≈ 2.0) ──────────────────────────────────────────────
    "ecotox_dl"   : ( 2.0,   0.0,  4.4, 1.2),
    "sqlite"      : ( 2.0,  -2.1,  3.6, 1.9),
    "install_r"   : ( 2.0,  -4.5,  4.4, 1.3),
    "units"       : ( 2.0,  -7.0,  4.6, 2.6),
    "calc_medians": ( 2.0,  -9.9,  4.4, 1.2),
    "fst_ecotox"  : ( 2.0, -11.8,  3.6, 1.9),

    # ── Online chain (x ≈ 10–15) ──────────────────────────────────────────────
    "user_input"  : (10.5,   0.0,  4.4, 1.3),
    "orient_dec"  : (10.5,  -2.0,  3.2, 1.5),
    "transpose"   : (14.8,  -2.0,  3.6, 1.4),
    "tidy"        : (10.5,  -3.9,  4.0, 1.2),
    "layer1"      : (10.5,  -5.6,  4.2, 1.6),
    "known_cas"   : (16.4,  -5.6,  3.4, 1.9),
    "resolved1"   : (10.5,  -7.8,  3.2, 1.5),
    "layer2"      : (15.1,  -7.8,  4.0, 1.8),
    "resolved2"   : (15.1, -10.2,  3.2, 1.5),
    "merge_cas"   : (15.1, -11.9,  4.2, 1.2),
    "manual"      : (20.0, -11.9,  3.8, 1.4),
    "join_ecotox" : (10.5, -12.5,  4.6, 1.4),

    # ── Taxonomic blocks ──────────────────────────────────────────────────────
    "tu_algae"    : ( 5.5, -17.8,  4.0, 2.2),
    "pti_algae"   : ( 5.5, -20.8,  4.0, 1.3),
    "tu_crust"    : (10.5, -17.8,  4.0, 2.2),
    "pti_crust"   : (10.5, -20.8,  4.0, 1.3),
    "tu_fish"     : (15.5, -17.8,  4.0, 2.2),
    "pti_fish"    : (15.5, -20.8,  4.0, 1.3),

    # ── Output layer ──────────────────────────────────────────────────────────
    "plots"       : ( 5.0, -23.6,  4.2, 1.9),
    "tables"      : (10.5, -23.6,  4.4, 1.9),
    "excel"       : (16.0, -23.6,  4.6, 1.9),
}

def nd(name):
    """Return (tx, ty, w, h) for node *name*."""
    return N[name]

def cx(name): return N[name][0]
def cy(name): return N[name][1]
def nw(name): return N[name][2]
def nh(name): return N[name][3]

# Edge shortcuts (TikZ coords of node borders)
def top(name):    return cy(name) + nh(name) / 2
def bot(name):    return cy(name) - nh(name) / 2
def lft(name):    return cx(name) - nw(name) / 2
def rgt(name):    return cx(name) + nw(name) / 2


# =============================================================================
# MAIN
# =============================================================================

def main():
    prs = Presentation()

    # A3 landscape
    prs.slide_width  = _cm(W_SLIDE)
    prs.slide_height = _cm(H_SLIDE)

    blank_layout = prs.slide_layouts[6]   # completely blank
    slide = prs.slides.add_slide(blank_layout)

    # ── Title ─────────────────────────────────────────────────────────────────
    tb = slide.shapes.add_textbox(_cm(1.0), _cm(0.3), _cm(35), _cm(1.0))
    _no_line(tb)
    tf = tb.text_frame
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "TaxoTox \u2014 Architectural Workflow"
    run.font.size  = Pt(16)
    run.font.bold  = True
    run.font.name  = "Calibri"
    run.font.color.rgb = BLACK

    tb2 = slide.shapes.add_textbox(_cm(1.0), _cm(1.15), _cm(38), _cm(0.7))
    _no_line(tb2)
    p2  = tb2.text_frame.paragraphs[0]
    run2 = p2.add_run()
    run2.text = ("Mixture Toxicity Assessment of Aquatic Pollutant Assemblages "
                 "using the Toxic Unit (TU) Framework")
    run2.font.size   = Pt(9)
    run2.font.italic = True
    run2.font.name   = "Calibri"
    run2.font.color.rgb = MGRAY

    # =========================================================================
    # BACKGROUND CONTAINERS  (drawn first — behind all nodes)
    # =========================================================================

    # Offline box: x −0.6 → 4.6, y 2.5 → −12.85
    add_box(slide, -0.6, 2.5, 4.6, -12.85,
            fill=OL_F, border=OL_B, border_pt=1.6, fill_alpha=35)

    # Online box: x 7.8 → 22.6, y 1.5 → −13.6
    add_box(slide, 7.8, 1.5, 22.6, -13.6,
            fill=ON_F, border=ON_B, border_pt=1.6, fill_alpha=25)

    # Taxonomic box: x 3.0 → 18.0, y −16.6 → −22.0
    add_box(slide, 3.0, -16.6, 18.0, -22.0,
            fill=PG_F, border=PG_B, border_pt=1.4, dashed=True, fill_alpha=18)

    # Algae / Crust / Fish sub-boxes
    for xc, fc, bc in ((5.5, AL_F, AL_B), (10.5, CR_F, CR_B), (15.5, FI_F, FI_B)):
        add_box(slide,
                xc - 2.3, top("tu_algae") + 0.2,
                xc + 2.3, bot("pti_algae") - 0.2,
                fill=fc, border=bc, border_pt=1.2, fill_alpha=30)

    # Output box: x 2.6 → 18.8, y −22.65 → −25.55
    add_box(slide, 2.6, -22.65, 18.8, -25.55,
            fill=OU_F, border=OU_B, border_pt=1.4, fill_alpha=25)

    # ── Section labels (placed after boxes so they sit on top) ────────────────
    add_label(slide, -0.55, 2.4,
              "OFFLINE PREPARATION PHASE\n(run once, refresh quarterly)",
              size=7.5, bold=True, color=OL_B)
    add_label(slide, 7.85, 1.4,
              "ONLINE SESSION PHASE   (interactive Shiny app)",
              size=7.5, bold=True, color=ON_B)
    add_label(slide, 3.05, -16.65,
              "PARALLEL TAXONOMIC COMPUTATION",
              size=7.5, bold=True, color=PG_B)
    add_label(slide, 2.65, -22.68,
              "OUTPUT LAYER",
              size=7.5, bold=True, color=OU_B)

    # =========================================================================
    # NODES
    # =========================================================================

    # ── Offline phase ─────────────────────────────────────────────────────────
    add_proc(slide, *nd("ecotox_dl"), PG_F, PG_B, [
        "Use ecotoxR to update",
        "SQLite EPA/ECOTOX database",
    ])

    add_db(slide, *nd("sqlite"), [
        "ECOTOX SQLite",
        "EPA Aquatic",
        "LC\u2085\u2080 / EC\u2085\u2080",
    ])

    add_proc(slide, *nd("install_r"), PG_F, PG_B, [
        "taxotox_install.R",
        "SQL extraction",
        "(fish, algae, crustacean)",
    ])

    add_proc(slide, *nd("units"), PG_F, PG_B, [
        "Unit Normalisation",
        "ng L\u207b\u00b9   \u03bcg L\u207b\u00b9   mg L\u207b\u00b9",
        [("\u2198  \u2193  \u2199", False, False, None)],
        "ng L\u207b\u00b9",
        [("back-transform (log)LC\u2085\u2080", False, True, None)],
    ])

    add_proc(slide, *nd("calc_medians"), PG_F, PG_B, [
        "Calculate Medians",
        "per taxon and pollutant",
    ])

    add_db(slide, *nd("fst_ecotox"), [
        "final_ecotox_data.fst",
        "conc_ng_L per",
        "(cas_number, ecotox_group)",
    ])

    # ── Online phase ──────────────────────────────────────────────────────────
    add_para(slide, *nd("user_input"), [
        "User Uploads File",
        "CSV / XLSX",
        "Units: ng L\u207b\u00b9",
    ])

    add_diamond(slide, *nd("orient_dec"), [
        [("Data\norientation?", False, False, None)],
    ])

    add_proc(slide, *nd("transpose"), PG_F, PG_B, [
        "Transpose matrix",
        "(row layout \u2192",
        "canonical layout)",
    ])

    add_proc(slide, *nd("tidy"), PG_F, PG_B, [
        "Tidy & Validate",
        "Parse numeric values",
    ])

    add_proc(slide, *nd("layer1"), PG_F, PG_B, [
        "CASRN Layer 1",
        "Exact match in",
        "Known_CAS.fst",
        [("Status: Confirmed", False, True, None)],
    ])

    add_db(slide, *nd("known_cas"), [
        "Known_CAS.fst",
        "Curated internal",
        "CASRN lookup",
    ])

    add_diamond(slide, *nd("resolved1"), [
        [("All\nresolved?", False, False, None)],
    ])

    add_proc(slide, *nd("layer2"), PG_F, PG_B, [
        "Layer 2 (Optional)",
        "PubChem REST API",
        "via webchem",
        "(Web search checkbox)",
    ])

    add_diamond(slide, *nd("resolved2"), [
        [("All\nresolved?", False, False, None)],
    ])

    add_proc(slide, *nd("merge_cas"), PG_F, PG_B, [
        "Merge Resolved CASRNs",
        "with concentration matrix",
    ])

    add_proc(slide, *nd("manual"), PG_F, PG_B, [
        "Manual CASRN Entry",
        "CASRN Matching tab",
        "User assigns CASRNs",
    ])

    add_proc(slide, *nd("join_ecotox"), PG_F, PG_B, [
        "Join with ECOTOX data",
        "Match on cas_number",
        "Group: algae / crustacean / fish",
    ])

    # ── Taxonomic blocks ──────────────────────────────────────────────────────
    TU_FORMULA = "TU\u1d62\u2c7c = C\u1d62\u2c7c / LC\u0303\u2085\u2080,\u1d62"
    PTI_FORMULA = "PTI\u2c7c = \u03a3\u1d62 TU\u1d62\u2c7c"

    add_proc(slide, *nd("tu_algae"), AL_F, AL_B, [
        [("Algae TU", True, False, AL_D)],
        TU_FORMULA,
        "per compound, per station",
    ])
    add_proc(slide, *nd("pti_algae"), AL_F, AL_B, [
        [("Algae PTI", True, False, AL_D)],
        PTI_FORMULA,
    ])

    add_proc(slide, *nd("tu_crust"), CR_F, CR_B, [
        [("Crustacean TU", True, False, CR_D)],
        TU_FORMULA,
        "per compound, per station",
    ])
    add_proc(slide, *nd("pti_crust"), CR_F, CR_B, [
        [("Crustacean PTI", True, False, CR_D)],
        PTI_FORMULA,
    ])

    add_proc(slide, *nd("tu_fish"), FI_F, FI_B, [
        [("Fish TU", True, False, FI_D)],
        TU_FORMULA,
        "per compound, per station",
    ])
    add_proc(slide, *nd("pti_fish"), FI_F, FI_B, [
        [("Fish PTI", True, False, FI_D)],
        PTI_FORMULA,
    ])

    # ── Output layer ──────────────────────────────────────────────────────────
    add_proc(slide, *nd("plots"), OU_F, OU_B, [
        [("Interactive Plots", True, False, RR)],
        "ggplot2 bar charts",
        "Per-sample PTI",
        "Per-pollutant TU",
    ])

    add_proc(slide, *nd("tables"), OU_F, OU_B, [
        [("DT Risk Tables", True, False, RR)],
        [("PTI \u2265 1.0: ", False, False, None), ("Acute risk", True, False, RR)],
        [("PTI 0.1\u20131.0: ", False, False, None), ("Chronic risk", True, False, RA)],
        [("PTI < 0.1: ", False, False, None), ("Low risk", True, False, RG)],
    ])

    add_proc(slide, *nd("excel"), OU_F, OU_B, [
        [("Excel Export (.xlsx)", True, False, RR)],
        "Legend / Original data",
        "Algae / Crust. / Fish",
        "CASRN Report",
        "Coverage Matrix",
    ])

    # =========================================================================
    # ARROWS
    # =========================================================================

    # ── Offline chain ─────────────────────────────────────────────────────────
    arrow(slide, cx("ecotox_dl"),    bot("ecotox_dl"),
                 cx("sqlite"),       top("sqlite"),       OL_B)
    arrow(slide, cx("sqlite"),       bot("sqlite"),
                 cx("install_r"),    top("install_r"),    OL_B)
    arrow(slide, cx("install_r"),    bot("install_r"),
                 cx("units"),        top("units"),        OL_B)
    arrow(slide, cx("units"),        bot("units"),
                 cx("calc_medians"), top("calc_medians"), OL_B)
    arrow(slide, cx("calc_medians"), bot("calc_medians"),
                 cx("fst_ecotox"),   top("fst_ecotox"),   OL_B)

    # fst_ecotox → join_ecotox  (dashed, horizontal then down)
    arrow(slide, rgt("fst_ecotox"), cy("fst_ecotox"),
                 lft("join_ecotox"), cy("join_ecotox"),   PG_B, dashed=True)

    # known_cas → layer1  (dashed)
    arrow(slide, lft("known_cas"), cy("known_cas"),
                 rgt("layer1"),    cy("layer1"),          PG_B, dashed=True)

    # ── Online chain ──────────────────────────────────────────────────────────
    arrow(slide, cx("user_input"),  bot("user_input"),
                 cx("orient_dec"), top("orient_dec"),    ON_B)

    # orient_dec → tidy (col layout, straight down)
    arrow(slide, cx("orient_dec"), bot("orient_dec"),
                 cx("tidy"),       top("tidy"),          ON_B)
    arrow_label(slide, cx("orient_dec") + 0.15, (bot("orient_dec") + top("tidy")) / 2,
                "col layout")

    # orient_dec → transpose (row layout, right)
    arrow(slide, rgt("orient_dec"),  cy("orient_dec"),
                 lft("transpose"),   cy("transpose"),     ON_B)
    arrow_label(slide, (rgt("orient_dec") + lft("transpose")) / 2 - 0.9,
                cy("orient_dec") + 0.2, "row layout")

    # transpose → tidy (down then left)
    arrow(slide, cx("transpose"), bot("transpose"),
                 rgt("tidy"),      cy("tidy"),            ON_B)

    arrow(slide, cx("tidy"),   bot("tidy"),
                 cx("layer1"), top("layer1"),             ON_B)
    arrow(slide, cx("layer1"),   bot("layer1"),
                 cx("resolved1"), top("resolved1"),       ON_B)

    # resolved1 → layer2  (No, right)
    arrow(slide, rgt("resolved1"), cy("resolved1"),
                 lft("layer2"),    cy("layer2"),          ON_B)
    arrow_label(slide, (rgt("resolved1") + lft("layer2")) / 2 - 0.9,
                cy("resolved1") + 0.2, "No")

    # resolved1 → merge_cas (Yes, diagonal)
    arrow(slide, cx("resolved1"), bot("resolved1"),
                 cx("merge_cas"), top("merge_cas"),       ON_B)
    arrow_label(slide, cx("resolved1") + 0.15,
                (bot("resolved1") + top("merge_cas")) / 2, "Yes")

    arrow(slide, cx("layer2"),   bot("layer2"),
                 cx("resolved2"), top("resolved2"),       ON_B)

    # resolved2 → merge_cas  (Yes, straight down)
    arrow(slide, cx("resolved2"), bot("resolved2"),
                 cx("merge_cas"), top("merge_cas"),       ON_B)
    arrow_label(slide, cx("resolved2") + 0.15,
                (bot("resolved2") + top("merge_cas")) / 2, "Yes")

    # resolved2 → manual  (No, right)
    arrow(slide, rgt("resolved2"), cy("resolved2"),
                 lft("manual"),    cy("manual"),          ON_B)
    arrow_label(slide, (rgt("resolved2") + lft("manual")) / 2 - 0.9,
                cy("resolved2") + 0.2, "No")

    # manual → merge_cas  (left)
    arrow(slide, lft("manual"),    cy("manual"),
                 rgt("merge_cas"), cy("merge_cas"),       ON_B)

    # merge_cas → join_ecotox  (left and down)
    arrow(slide, cx("merge_cas"),  bot("merge_cas"),
                 cx("join_ecotox"), top("join_ecotox"),   ON_B)

    # ── join_ecotox → taxonomic blocks ────────────────────────────────────────
    arrow(slide, cx("join_ecotox"), bot("join_ecotox"),
                 cx("tu_algae"),    top("tu_algae"),       AL_B)
    arrow(slide, cx("join_ecotox"), bot("join_ecotox"),
                 cx("tu_crust"),    top("tu_crust"),       CR_B)
    arrow(slide, cx("join_ecotox"), bot("join_ecotox"),
                 cx("tu_fish"),     top("tu_fish"),        FI_B)

    # ── TU → PTI ──────────────────────────────────────────────────────────────
    arrow(slide, cx("tu_algae"),  bot("tu_algae"),
                 cx("pti_algae"), top("pti_algae"),        AL_B)
    arrow(slide, cx("tu_crust"),  bot("tu_crust"),
                 cx("pti_crust"), top("pti_crust"),        CR_B)
    arrow(slide, cx("tu_fish"),   bot("tu_fish"),
                 cx("pti_fish"),  top("pti_fish"),         FI_B)

    # ── PTI → output box (bottom of output container) ─────────────────────────
    out_box_top = -22.65   # top edge of output box (TikZ y)
    arrow(slide, cx("pti_algae"), bot("pti_algae"),
                 cx("pti_algae"), out_box_top,             OU_B)
    arrow(slide, cx("pti_crust"), bot("pti_crust"),
                 cx("pti_crust"), out_box_top,             OU_B)
    arrow(slide, cx("pti_fish"),  bot("pti_fish"),
                 cx("pti_fish"),  out_box_top,             OU_B)

    # =========================================================================
    # INFO BOXES  (formula + legend — to the right of the main diagram)
    # =========================================================================

    add_textbox(slide, 23.3, -4.5, 5.8, 5.8,
        paragraphs=[
            [("Key Equations:", True, False, None)],
            "",
            [("TU\u1d62\u2c7c  =  C\u1d62\u2c7c  /  LC\u0303\u2085\u2080,\u1d62", False, False, None)],
            "",
            [("C\u1d62\u2c7c", False, False, None),
             (" = conc. of compound i, sample j (ng L\u207b\u00b9)", False, False, None)],
            [("LC\u0303\u2085\u2080,\u1d62", False, False, None),
             (" = median(LC\u2085\u2080/EC\u2085\u2080)", False, False, None)],
            "   compound i, per taxon group",
            "",
            [("PTI\u2c7c  =  \u03a3\u1d62  TU\u1d62\u2c7c", False, False, None)],
            "",
            [("Risk thresholds:", True, False, None)],
            [("\u2022 PTI \u2265 1.0: Acute risk", False, False, RR)],
            [("\u2022 PTI 0.1\u20131.0: Chronic risk", False, False, RA)],
            [("\u2022 PTI < 0.1: Low risk", False, False, RG)],
        ],
        base_size=7.5, align=PP_ALIGN.LEFT,
    )

    add_textbox(slide, 23.3, -13.5, 5.8, 4.5,
        paragraphs=[
            [("Symbol Legend:", True, False, None)],
            "",
            "Parallelogram: data input",
            "Rectangle: processing step",
            "Cylinder: database / file store",
            "Diamond: decision point",
            "",
            [("Taxonomic colour coding:", True, False, None)],
            [("\u25a0  Algae", False, False, AL_D)],
            [("\u25a0  Crustaceans", False, False, CR_D)],
            [("\u25a0  Fish", False, False, FI_D)],
        ],
        base_size=7.5, align=PP_ALIGN.LEFT,
    )

    # ── Attribution ───────────────────────────────────────────────────────────
    tb = slide.shapes.add_textbox(_cm(1.0), _cm(H_SLIDE - 0.55), _cm(38), _cm(0.45))
    _no_line(tb)
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = ("Gridish N. & Suari Y.  \u2014  School of Marine Sciences, "
                "Ruppin Academic Center, Israel  \u2014  MIT License")
    run.font.size  = Pt(6.5)
    run.font.name  = "Calibri"
    run.font.color.rgb = MGRAY

    # ── Save ──────────────────────────────────────────────────────────────────
    prs.save(str(OUT))
    print(f"Saved: {OUT}")


if __name__ == "__main__":
    main()
